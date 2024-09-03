import tkinter as tk
from tkinter import messagebox
import os
import shutil
from openpyxl import load_workbook
from openpyxl.styles import PatternFill
import pandas as pd
from pptx import Presentation
from pptx.util import Pt, Cm, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import qrcode
import time
from datetime import datetime
import io

# Obter ano e mês atuais
current_year = datetime.now().strftime('%Y')
current_month = datetime.now().strftime('%m')

# Construir diretório de origem com base no ano e mês atuais
source_directory = rf'\\caminho\{current_year}\{current_month}-{current_year}'
ppt_directory = r'\\caminho\CRACHAS'
#destination_directory = os.path.join(os.getenv('TEMP'), 'crachas')  # Usar diretório temporário do sistema
destination_directory = rf'\\caminho\CRACHAS\CRACHÁS FEITOS\{current_year}\{current_month}-{current_year}\CRACHAS AUTOMATICOS'
qr_code_directory = os.path.join(destination_directory, 'qr_code')
ready_directory = os.path.join(destination_directory, 'prontos')

if not os.path.exists(destination_directory):
    os.makedirs(destination_directory)

# Função para iniciar o processamento dos crachás
def iniciar_crachas():
    os.makedirs(qr_code_directory, exist_ok=True)
    os.makedirs(ready_directory, exist_ok=True)
    
    def copy_files(source_folder, destination_folder):
        for file_name in os.listdir(source_folder):
            if file_name.startswith('CRACHÁS') and file_name.endswith('.xlsx'):
                source_file = os.path.join(source_folder, file_name)
                destination_file = os.path.join(destination_folder, file_name)
                time.sleep(1)
                if os.path.exists(destination_file):
                    print(f'{destination_file} já existe')
                else:
                    shutil.copy2(source_file, destination_file)
                    print(f'Arquivo copiado: {source_file} para {destination_file}')

    def remove_filters(workbook_path):
        wb = load_workbook(workbook_path)
        for sheet in wb.worksheets:
            if sheet.auto_filter:
                sheet.auto_filter.ref = None
        wb.save(workbook_path)

    def generate_qr_code(matricula):
        qr = qrcode.QRCode(version=1, box_size=10, border=5)
        qr.add_data(matricula)
        qr.make(fit=True)
        img = qr.make_image(fill='black', back_color='white')
        qr_code_path = os.path.join(qr_code_directory, f"{matricula}.png")
        img.save(qr_code_path)
        return qr_code_path

    def truncate_function(text, max_length):
        words = text.split()
        if len(text) > max_length:
            if len(words) > 1:
                words[1] = words[1][:3] + '.'
            if len(' '.join(words)) > max_length and len(words) > 3:
                words[3] = words[3][:3] + '.'
            return ' '.join(words)
        return text

    def truncate_name(text, max_length):
        words = text.split()
        if len(text) > max_length:
            if len(words) > 1:
                words[1] = words[1][:1] + '.'
            if len(' '.join(words)) > max_length and len(words) > 2:
                words[2] = words[2][:1] + '.'
            return ' '.join(words)
        return text


    def update_ppt(prs, ppt_template_path, data_row, row_index, sheet):
        
        # Carregar a apresentação do template
        template_prs = Presentation(ppt_template_path)

        # Definir as dimensões exatas do slide (em centímetros)
        slide_width = Cm(5.499)
        slide_height = Cm(8.502)
        prs.slide_width = slide_width
        prs.slide_height = slide_height
        template_prs.slide_width = slide_width
        template_prs.slide_height = slide_height

        matricula_formatada = str(data_row['Matricula']).zfill(6)
        qr_code_path = generate_qr_code(matricula_formatada)

        for slide in template_prs.slides:
            # Criar um novo slide em branco na apresentação principal
            new_slide = prs.slides.add_slide(prs.slide_layouts[5])
            first_slide = False

            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Criar uma nova caixa de texto com o mesmo tamanho e posição
                    new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                    new_text_frame = new_shape.text_frame
                    new_text_frame.word_wrap = shape.text_frame.word_wrap
                    new_shape.top = (shape.top -Pt(20))
                    for paragraph in shape.text_frame.paragraphs:
                        new_paragraph = new_text_frame.add_paragraph()
                        new_paragraph.alignment = paragraph.alignment

                        for run in paragraph.runs:
                            new_run = new_paragraph.add_run()
                            new_run.text = run.text
                            new_run.font.name = run.font.name
                            new_run.font.size = run.font.size
                            new_run.font.bold = run.font.bold

                            try:
                                if run.font.color and run.font.color.rgb:
                                    new_run.font.color.rgb = run.font.color.rgb
                            except AttributeError:
                                new_run.font.color.rgb = RGBColor(0x4F, 0x1E, 0x0B)

                        # Realizar substituições de texto
                        if 'MATCMP' in new_paragraph.text:
                            first_slide = new_slide
                            new_paragraph.text = new_paragraph.text.replace('MATCMP', matricula_formatada)
                            # Configuração da fonte para matrícula
                            for run in new_paragraph.runs:
                                run.font.name = 'Calibri'
                                run.font.size = Pt(12)
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(0x4F, 0x1E, 0x0B)

                        if 'NOME COMPLETO' in new_paragraph.text:
                            nome_completo = truncate_name(data_row['Nome complet'], 30)
                            new_paragraph.text = new_paragraph.text.replace('NOME COMPLETO', nome_completo)
                            new_paragraph.alignment = PP_ALIGN.CENTER
                            # Configuração da fonte para nome completo
                            for run in new_paragraph.runs:
                                run.font.name = 'Calibri'
                                run.font.size = Pt(8)
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(0x4F, 0x1E, 0x0B)

                        if 'CARGO COMPLETO' in new_paragraph.text:
                            cargo = truncate_function(data_row['Desc. Comple'].replace(' III', '').replace(' II', '').replace(' I', ''), 30)
                            new_paragraph.text = new_paragraph.text.replace('CARGO COMPLETO', cargo)
                            new_paragraph.alignment = PP_ALIGN.CENTER
                            # Configuração da fonte para cargo completo
                            for run in new_paragraph.runs:
                                run.font.name = 'Calibri'
                                run.font.size = Pt(7)
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(0x4F, 0x1E, 0x0B)
                        

                elif shape.shape_type == 13:  # Tipo de forma: Imagem
                    # Copiar as imagens para o novo slide
                    image_stream = io.BytesIO(shape.image.blob)
                    new_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)

            # Adicionar o QR code ao slide
            if first_slide and os.path.exists(qr_code_path):
                try:
                    left = Cm(0.49)
                    top = Cm(6.5)
                    width = Cm(1.72)
                    height = Cm(1.72)
                    first_slide.shapes.add_picture(qr_code_path, left, top, width, height)
                    mark_row_as_processed(sheet, row_index)
                except Exception as e:
                    print(f"Erro ao adicionar o QR code: {e}")
            else:
                mark_row_as_processed(sheet, row_index)



    def mark_row_as_processed(sheet, row_index):
        fill = PatternFill(start_color='00FF00', end_color='00FF00', fill_type='solid')
        for cell in sheet[row_index]:
            cell.fill = fill

    def map_columns(df):
        df.dropna(how='all', inplace=True)
        df.reset_index(drop=True, inplace=True)
        columns_map = {
            'Matricula': 'Matricula',
            'Nome complet': 'Nome complet',
            'Desc. Comple': 'Desc. Comple'
        }
        for i, row in df.iterrows():
            if 'Matricula' in row.values and 'Nome complet' in row.values and 'Desc. Comple' in row.values:
                df.columns = row
                df = df[i + 1:]
                df.reset_index(drop=True, inplace=True)
                print(f"Colunas mapeadas corretamente: {list(df.columns)}")
                return df

        print(f"Erro ao mapear colunas. Dados encontrados: {df.head()}")
        raise ValueError('As colunas necessárias não foram encontradas.')

    def process_spreadsheets(spreadsheet_path):
        wb = load_workbook(spreadsheet_path)
        sheet = wb.worksheets[0]
        df = pd.DataFrame(sheet.values)
        try:
            df = map_columns(df)
            prs = Presentation()  # Criar uma apresentação PowerPoint para cada planilha Excel
            for index, row in df.iterrows():
                ppt_template_path = os.path.join(ppt_directory, "PADRAO CRACHA.pptx")
                update_ppt(prs, ppt_template_path, row, index + 2, sheet)
            
            # Salvar a apresentação após processar todas as linhas da planilha atual
            output_path = os.path.join(ready_directory, f"{os.path.splitext(os.path.basename(spreadsheet_path))[0]}.pptx")
            prs.save(output_path)
            print(f"Arquivo PowerPoint gerado para a planilha {spreadsheet_path}: {output_path}")
            wb.save(spreadsheet_path)
        except Exception as e:
            print(f"Erro ao processar a planilha {spreadsheet_path}: {e}")


    copy_files(source_directory, destination_directory)
    for file_name in os.listdir(destination_directory):
        if file_name.startswith('CRACHÁS') and file_name.endswith('.xlsx'):
            spreadsheet_path = os.path.join(destination_directory, file_name)
            remove_filters(spreadsheet_path)
            process_spreadsheets(spreadsheet_path)

    messagebox.showinfo("Sucesso", f"Processamento de crachás concluído na pasta {ready_directory}!")
# Criação da interface gráfica
root = tk.Tk()
root.title("Gerador de Crachás")

button = tk.Button(root, text="INICIAR CRACHÁS", command=iniciar_crachas, height=2, width=20)
button.pack(pady=20)

root.mainloop()
